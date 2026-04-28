from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\sseja\OneDrive\Documents\New project 2")
TEMPLATE = Path(r"C:\Users\sseja\Downloads\[EXT] Solution Challenge 2026 - Prototype PPT Template (1).pptx")
OUTPUT = ROOT / "docs" / "SecureMedia_AI_Solution_Challenge_Deck.pptx"
ASSETS = ROOT / "docs" / "assets"
SERVICE_URL = "https://securemedia-ai-97902534410.us-central1.run.app"

FONT = "Google Sans"
MONO = "Consolas"
BLUE = RGBColor(66, 133, 244)
GREEN = RGBColor(52, 168, 83)
YELLOW = RGBColor(251, 188, 4)
RED = RGBColor(234, 67, 53)
DARK = RGBColor(32, 33, 36)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(248, 250, 253)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(218, 220, 224)
NAVY = RGBColor(27, 45, 72)


def set_run_font(run, size, color=DARK, bold=False, italic=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def set_title(shape, text):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, 18, color=DARK, bold=False)


def clear_shape_text(shape):
    if not hasattr(shape, "text_frame"):
        return
    shape.text_frame.clear()


def add_text_lines(slide, x, y, w, h, lines, font=FONT, fill=None, line=None, radius=True):
    if fill is not None:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            x,
            y,
            w,
            h,
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        if line is None:
            box.line.fill.background()
        else:
            box.line.color.rgb = line
            box.line.width = Pt(1)
        frame = box.text_frame
    else:
        frame = slide.shapes.add_textbox(x, y, w, h).text_frame

    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(6)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for idx, item in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(item.get("space_after", 4))
        paragraph.space_before = Pt(item.get("space_before", 0))
        if item.get("bullet"):
            paragraph.level = item.get("level", 0)
            paragraph.bullet = True
        run = paragraph.add_run()
        run.text = item["text"]
        set_run_font(
            run,
            item.get("size", 16),
            color=item.get("color", DARK),
            bold=item.get("bold", False),
            italic=item.get("italic", False),
            font=item.get("font", font),
        )
        if item.get("url"):
            run.hyperlink.address = item["url"]

    return frame


def add_card(slide, x, y, w, h, title, bullets, fill):
    lines = [{"text": title, "size": 14, "bold": True, "color": NAVY, "space_after": 6}]
    for bullet in bullets:
        lines.append({"text": bullet, "size": 11.5, "color": DARK, "bullet": True, "space_after": 2})
    add_text_lines(slide, x, y, w, h, lines, fill=fill, line=BORDER)


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
    return arrow


def add_picture_contain(slide, image_path, x, y, w, h, border=False):
    image = Image.open(image_path)
    img_w, img_h = image.size
    box_ratio = w / h
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        new_w = w
        new_h = int(w / img_ratio)
        new_x = x
        new_y = y + int((h - new_h) / 2)
    else:
        new_h = h
        new_w = int(h * img_ratio)
        new_y = y
        new_x = x + int((w - new_w) / 2)

    if border:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER
        box.line.width = Pt(1)
    slide.shapes.add_picture(str(image_path), new_x, new_y, width=new_w, height=new_h)


def add_flow_box(slide, x, y, w, h, title, subtitle, fill):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    p1 = frame.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_font(r1, 12, color=NAVY, bold=True)
    p2 = frame.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = subtitle
    set_run_font(r2, 9.5, color=DARK)


def add_table_row(slide, x, y, widths, height, cells, fills):
    current_x = x
    for idx, cell in enumerate(cells):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, current_x, y, widths[idx], height)
        box.fill.solid()
        box.fill.fore_color.rgb = fills[idx]
        box.line.color.rgb = BORDER
        box.line.width = Pt(1)
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Pt(6)
        frame.margin_right = Pt(6)
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if idx < 2 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = cell
        set_run_font(r, 11 if idx < 2 else 12, color=DARK, bold=(idx == 0))
        current_x += widths[idx] + Pt(6)


def build():
    prs = Presentation(str(TEMPLATE))

    # Slide 1: keep template guidelines.

    # Slide 2
    slide = prs.slides[1]
    clear_shape_text(slide.shapes[2])
    add_text_lines(
        slide,
        Inches(0.35),
        Inches(3.08),
        Inches(9.15),
        Inches(1.95),
        [
            {"text": "Team Details", "size": 18, "bold": True, "color": DARK, "space_after": 10},
            {"text": "Team name: double trouble 2x", "size": 13, "color": NAVY, "bold": True, "space_after": 4},
            {"text": "Team members:", "size": 12.5, "color": DARK, "bold": True, "space_after": 2},
            {"text": "1) Aadi shah", "size": 12.5, "color": DARK, "space_after": 2},
            {"text": "2) nitya singh", "size": 12.5, "color": DARK, "space_after": 4},
            {
                "text": (
                    "Problem Statement: creators, journalists, marketplaces, and moderators "
                    "need a fast way to detect duplicate media and verify ownership before "
                    "content is reused, sold, or published."
                ),
                "size": 12.5,
                "color": DARK,
                "space_after": 2,
            },
        ],
    )

    # Slide 3
    slide = prs.slides[2]
    set_title(slide.shapes[0], "Brief about SecureMedia AI")
    add_text_lines(
        slide,
        Inches(0.5),
        Inches(1.65),
        Inches(5.7),
        Inches(2.5),
        [
            {
                "text": (
                    "SecureMedia AI is a lightweight authenticity workflow for image uploads. "
                    "The product combines perceptual hashing, optional AI similarity, and "
                    "blockchain ownership status in one review step."
                ),
                "size": 17,
                "color": DARK,
                "space_after": 8,
            },
            {
                "text": (
                    "The MVP now runs as a live Google Cloud deployment with a React upload "
                    "interface, a Flask upload API, duplicate detection, ownership verification "
                    "fields, and a minimal smart contract path for EVM testnets."
                ),
                "size": 15,
                "color": GRAY,
                "space_after": 8,
            },
            {
                "text": (
                    "For Google Solution Challenge alignment, the production path now uses "
                    "Cloud Run hosting and is ready for Vertex AI multimodal similarity as the "
                    "Google-native scoring layer."
                ),
                "size": 14,
                "color": BLUE,
            },
        ],
    )
    add_card(slide, Inches(6.45), Inches(1.75), Inches(1.0), Inches(1.1), "Trust", ["Hash + AI + ownership"], LIGHT)
    add_card(slide, Inches(7.6), Inches(1.75), Inches(1.3), Inches(1.1), "Live", ["Cloud Run deployment"], LIGHT)
    add_card(slide, Inches(6.45), Inches(3.0), Inches(2.45), Inches(1.15), "Next Step", ["Vertex AI scoring", "Managed secrets"], LIGHT)

    # Slide 4
    slide = prs.slides[3]
    set_title(slide.shapes[0], "Opportunities and Differentiation")
    add_card(
        slide,
        Inches(0.45),
        Inches(1.7),
        Inches(3.0),
        Inches(2.7),
        "Different from similar ideas",
        [
            "Combines duplicate detection and ownership tracking in one workflow.",
            "Modular services keep the stack lightweight for pilots and hackathons.",
            "Still degrades gracefully when AI or blockchain modules are unavailable.",
        ],
        RGBColor(232, 240, 254),
    )
    add_card(
        slide,
        Inches(3.55),
        Inches(1.7),
        Inches(3.0),
        Inches(2.7),
        "How it solves the problem",
        [
            "Flags suspicious duplicate uploads before publication or resale.",
            "Adds ownership context to moderation and review decisions.",
            "Reduces ambiguity for creator-support, newsroom, and marketplace teams.",
        ],
        RGBColor(232, 245, 233),
    )
    add_card(
        slide,
        Inches(6.65),
        Inches(1.7),
        Inches(2.9),
        Inches(2.7),
        "USP",
        [
            "One API response with similarity, duplicate state, owner, and verification.",
            "Challenge-ready path to Google AI and Google Cloud.",
            "Simple enough to ship, strong enough to extend.",
        ],
        RGBColor(254, 247, 224),
    )

    # Slide 5
    slide = prs.slides[4]
    set_title(slide.shapes[0], "List of Features Offered by the Solution")
    feature_cards = [
        ("Image upload", "Drag-and-drop or picker-based image submission", RGBColor(232, 240, 254)),
        ("Live preview", "Visual confirmation before analysis begins", RGBColor(232, 245, 233)),
        ("Similarity score", "Clear percentage output for fast review", RGBColor(254, 247, 224)),
        ("Duplicate status", "Original vs duplicate decision signal", RGBColor(252, 232, 230)),
        ("Ownership status", "Owner string and blockchain verification field", RGBColor(232, 240, 254)),
        ("Modular backend", "Hashing, AI, and blockchain modules work independently", RGBColor(232, 245, 233)),
    ]
    positions = [
        (0.45, 1.7), (3.35, 1.7), (6.25, 1.7),
        (0.45, 3.15), (3.35, 3.15), (6.25, 3.15),
    ]
    for (title, body, fill), (x, y) in zip(feature_cards, positions):
        add_text_lines(
            slide,
            Inches(x),
            Inches(y),
            Inches(2.65),
            Inches(1.15),
            [
                {"text": title, "size": 13, "bold": True, "color": NAVY, "space_after": 6},
                {"text": body, "size": 11.2, "color": DARK},
            ],
            fill=fill,
            line=BORDER,
        )

    # Slide 6
    slide = prs.slides[5]
    set_title(slide.shapes[0], "Process Flow Diagram")
    boxes = [
        ("1. Upload", "User sends an image from the web client.", RGBColor(232, 240, 254)),
        ("2. Hash", "Perceptual hash is generated through the securemedia-core path.", RGBColor(232, 245, 233)),
        ("3. Compare", "Optional AI similarity compares against stored embeddings.", RGBColor(254, 247, 224)),
        ("4. Verify", "Ownership service checks and optionally registers the hash.", RGBColor(252, 232, 230)),
        ("5. Return", "Frontend receives similarity, duplicate, owner, verification.", RGBColor(232, 240, 254)),
    ]
    start_x = 0.35
    for idx, (title, subtitle, fill) in enumerate(boxes):
        add_flow_box(slide, Inches(start_x + idx * 1.9), Inches(2.1), Inches(1.65), Inches(1.1), title, subtitle, fill)
        if idx < len(boxes) - 1:
            add_arrow(slide, Inches(start_x + idx * 1.9 + 1.67), Inches(2.45), Inches(0.18), Inches(0.36))
    add_text_lines(
        slide,
        Inches(0.7),
        Inches(3.65),
        Inches(8.7),
        Inches(0.9),
        [
            {
                "text": "Primary users: creators, moderators, marketplaces, journalists, and authenticity review teams.",
                "size": 12.2,
                "color": GRAY,
                "align": PP_ALIGN.CENTER,
            }
        ],
    )

    # Slide 7
    slide = prs.slides[6]
    set_title(slide.shapes[0], "Wireframes / Mock Diagrams of the Proposed Solution")
    add_picture_contain(slide, ASSETS / "securemedia-ui-idle.png", Inches(0.45), Inches(1.6), Inches(5.8), Inches(3.35), border=True)
    add_text_lines(
        slide,
        Inches(6.5),
        Inches(1.75),
        Inches(2.9),
        Inches(0.85),
        [
            {"text": "Upload panel", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Single image input, preview area, and analyze trigger.", "size": 11.2, "color": DARK},
        ],
        fill=RGBColor(232, 240, 254),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(6.5),
        Inches(2.75),
        Inches(2.9),
        Inches(0.95),
        [
            {"text": "Decision cards", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Similarity, duplicate status, owner, and blockchain verification stay visible together.", "size": 11.2, "color": DARK},
        ],
        fill=RGBColor(232, 245, 233),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(6.5),
        Inches(3.9),
        Inches(2.9),
        Inches(0.9),
        [
            {"text": "Review posture", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Optimized for quick operator decisions, not consumer browsing.", "size": 11.2, "color": DARK},
        ],
        fill=RGBColor(254, 247, 224),
        line=BORDER,
    )

    # Slide 8
    slide = prs.slides[7]
    set_title(slide.shapes[0], "Architecture Diagram of the Proposed Solution")
    add_text_lines(slide, Inches(0.45), Inches(1.55), Inches(4.3), Inches(0.4), [{"text": "Current MVP", "size": 16, "bold": True, "color": NAVY}])
    add_text_lines(slide, Inches(5.15), Inches(1.55), Inches(4.25), Inches(0.4), [{"text": "Live Google Cloud Deployment", "size": 16, "bold": True, "color": NAVY}])
    mvp_layers = [
        (0.55, 2.0, 1.6, 0.75, "React Frontend", "Upload UI + results"),
        (2.35, 2.0, 1.7, 0.75, "Flask API", "POST /upload"),
        (0.55, 3.0, 1.6, 0.8, "Hash Service", "Perceptual hash + duplicate logic"),
        (2.35, 3.0, 1.7, 0.8, "AI Similarity", "Optional embedding comparison"),
        (1.45, 4.05, 1.8, 0.8, "Ownership Service", "Web3 + local fallback"),
    ]
    for x, y, w, h, title, subtitle in mvp_layers:
        add_flow_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), title, subtitle, RGBColor(232, 240, 254))
    add_arrow(slide, Inches(2.1), Inches(2.24), Inches(0.18), Inches(0.32))
    add_arrow(slide, Inches(2.1), Inches(3.24), Inches(0.18), Inches(0.32))
    add_arrow(slide, Inches(1.95), Inches(4.15), Inches(0.18), Inches(-0.48))
    add_text_lines(
        slide,
        Inches(5.2),
        Inches(2.0),
        Inches(1.8),
        Inches(0.8),
        [
            {"text": "Cloud Run", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Frontend + backend deployed together", "size": 10.8, "color": DARK},
        ],
        fill=RGBColor(232, 245, 233),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(7.1),
        Inches(2.0),
        Inches(1.8),
        Inches(0.8),
        [
            {"text": "Cloud Storage", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Planned evidence retention layer", "size": 10.8, "color": DARK},
        ],
        fill=RGBColor(254, 247, 224),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(5.2),
        Inches(3.05),
        Inches(1.8),
        Inches(0.8),
        [
            {"text": "Vertex AI Ready", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Google-native similarity service path already wired", "size": 10.8, "color": DARK},
        ],
        fill=RGBColor(252, 232, 230),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(7.1),
        Inches(3.05),
        Inches(1.8),
        Inches(0.8),
        [
            {"text": "Secrets / State", "size": 13, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "Secret Manager and persistence are the next hardening step", "size": 10.8, "color": DARK},
        ],
        fill=RGBColor(232, 240, 254),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(5.25),
        Inches(4.15),
        Inches(3.65),
        Inches(0.8),
        [
            {
                "text": "The current codebase already supports a modular split, which made the Google Cloud deployment a packaging step instead of a rewrite.",
                "size": 11.2,
                "color": GRAY,
                "align": PP_ALIGN.CENTER,
            }
        ],
        fill=LIGHT,
        line=BORDER,
    )

    # Slide 9
    slide = prs.slides[8]
    set_title(slide.shapes[0], "Technologies to be Used in the Solution")
    add_text_lines(slide, Inches(0.55), Inches(1.65), Inches(4.0), Inches(0.4), [{"text": "Current Prototype Stack", "size": 16, "bold": True, "color": NAVY}])
    add_text_lines(slide, Inches(5.2), Inches(1.65), Inches(4.0), Inches(0.4), [{"text": "Challenge Deployment Stack", "size": 16, "bold": True, "color": NAVY}])
    left_items = [
        "React 19 + Vite",
        "Tailwind CSS",
        "Flask API",
        "Pillow + ImageHash",
        "Optional Hugging Face CLIP",
        "Solidity + Web3",
        "Sepolia / Polygon testnets",
    ]
    right_items = [
        "Cloud Run (live)",
        "Cloud Storage",
        "Secret Manager",
        "Vertex AI multimodal embeddings",
        "Cloud Build / Artifact Registry",
        "Monitoring + Logging",
    ]
    for idx, item in enumerate(left_items):
        add_text_lines(
            slide,
            Inches(0.6 + (idx % 2) * 2.0),
            Inches(2.15 + (idx // 2) * 0.62),
            Inches(1.8),
            Inches(0.45),
            [{"text": item, "size": 11.5, "color": DARK, "align": PP_ALIGN.CENTER}],
            fill=RGBColor(232, 240, 254) if idx % 2 == 0 else RGBColor(232, 245, 233),
            line=BORDER,
        )
    for idx, item in enumerate(right_items):
        add_text_lines(
            slide,
            Inches(5.25 + (idx % 2) * 1.9),
            Inches(2.15 + (idx // 2) * 0.62),
            Inches(1.7),
            Inches(0.45),
            [{"text": item, "size": 11.5, "color": DARK, "align": PP_ALIGN.CENTER}],
            fill=RGBColor(254, 247, 224) if idx % 2 == 0 else RGBColor(252, 232, 230),
            line=BORDER,
        )

    # Slide 10
    slide = prs.slides[9]
    set_title(slide.shapes[0], "Estimated Implementation Cost")
    add_text_lines(
        slide,
        Inches(0.55),
        Inches(1.6),
        Inches(8.9),
        Inches(0.4),
        [{"text": "Estimated monthly cost for the live pilot on Google Cloud", "size": 14, "bold": True, "color": NAVY}],
    )
    widths = [Inches(2.3), Inches(4.35), Inches(1.55)]
    add_table_row(slide, Inches(0.55), Inches(2.05), widths, Inches(0.45), ["Component", "Assumption", "Monthly cost"], [LIGHT, LIGHT, LIGHT])
    cost_rows = [
        ("Cloud Run", "Live frontend + Flask API service in us-central1", "$10-$25"),
        ("Cloud Storage", "Uploaded image retention and evidence objects", "$2-$5"),
        ("Vertex AI", "Similarity calls when Google scoring is switched on", "$20-$40"),
        ("Observability", "Basic logs, traces, and monitoring", "$0-$10"),
        ("Testnet transactions", "Registration on Sepolia / Amoy", "Negligible"),
    ]
    for idx, row in enumerate(cost_rows):
        add_table_row(slide, Inches(0.55), Inches(2.55 + idx * 0.5), widths, Inches(0.42), list(row), [WHITE, WHITE, WHITE])
    add_text_lines(
        slide,
        Inches(0.55),
        Inches(5.0),
        Inches(8.9),
        Inches(0.45),
        [{"text": "Estimated pilot total: about $32-$80 per month, with hackathon credits likely reducing early cost.", "size": 12, "color": BLUE, "bold": True, "align": PP_ALIGN.CENTER}],
        fill=RGBColor(232, 240, 254),
        line=BORDER,
    )

    # Slide 11
    slide = prs.slides[10]
    set_title(slide.shapes[0], "Snapshots of the MVP")
    add_picture_contain(slide, ASSETS / "securemedia-cloudrun-live.png", Inches(0.45), Inches(1.6), Inches(4.45), Inches(2.55), border=True)
    add_picture_contain(slide, ASSETS / "securemedia-ui-file-selected.png", Inches(5.05), Inches(1.6), Inches(4.45), Inches(2.55), border=True)
    add_text_lines(
        slide,
        Inches(0.7),
        Inches(1.3),
        Inches(3.9),
        Inches(0.3),
        [{"text": "Live Cloud Run deployment", "size": 11.5, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
    )
    add_text_lines(
        slide,
        Inches(5.25),
        Inches(1.3),
        Inches(3.9),
        Inches(0.3),
        [{"text": "Local operator flow with file selected", "size": 11.5, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
    )
    add_text_lines(
        slide,
        Inches(0.65),
        Inches(4.25),
        Inches(8.7),
        Inches(0.7),
        [
            {"text": "Production smoke-test response", "size": 12.5, "bold": True, "color": NAVY, "space_after": 4},
            {"text": '{ "similarity": 0.0, "duplicate": false, "owner": "Unverified", "blockchain_verified": false }', "size": 11.5, "color": DARK, "font": MONO},
        ],
        fill=LIGHT,
        line=BORDER,
    )

    # Slide 12
    slide = prs.slides[11]
    set_title(slide.shapes[0], "Additional Details / Future Development")
    add_card(
        slide,
        Inches(0.55),
        Inches(1.75),
        Inches(4.1),
        Inches(2.9),
        "Immediate engineering follow-up",
        [
            "Turn on Vertex AI scoring in production with a managed embedding store.",
            "Replace placeholder chain config with a real contract deployment and funded wallet.",
            "Persist uploads, ownership history, and review trails beyond local JSON fallback.",
        ],
        RGBColor(232, 240, 254),
    )
    add_card(
        slide,
        Inches(4.9),
        Inches(1.75),
        Inches(4.1),
        Inches(2.9),
        "Challenge-ready roadmap",
        [
            "Add creator accounts, evidence history, and case management.",
            "Introduce Secret Manager and Cloud-native persistence for production operations.",
            "Extend from image uploads to video frames, watermarking, and newsroom workflows.",
        ],
        RGBColor(232, 245, 233),
    )

    # Slide 13
    slide = prs.slides[12]
    set_title(slide.shapes[0], "Project Links")
    link_rows = [
        ("GitHub Public Repository", "https://github.com/Aadishah17/secure-media", "https://github.com/Aadishah17/secure-media"),
        ("Demo Video Link (3 Minutes)", "To be recorded before submission", None),
        ("Live MVP Link", SERVICE_URL, SERVICE_URL),
        ("Health Check", f"{SERVICE_URL}/api/health", f"{SERVICE_URL}/api/health"),
    ]
    for idx, (label, value, url) in enumerate(link_rows):
        add_text_lines(
            slide,
            Inches(0.55),
            Inches(1.75 + idx * 0.78),
            Inches(8.9),
            Inches(0.62),
            [
                {"text": label, "size": 12.5, "bold": True, "color": NAVY, "space_after": 2},
                {"text": value, "size": 11.2, "color": BLUE if url else DARK, "url": url},
            ],
            fill=LIGHT if idx % 2 == 0 else WHITE,
            line=BORDER,
        )

    # Slide 14
    slide = prs.slides[13]
    add_text_lines(
        slide,
        Inches(0.55),
        Inches(0.95),
        Inches(3.7),
        Inches(0.55),
        [{"text": "Impact and SDG Alignment", "size": 24, "bold": True, "color": NAVY}],
    )
    add_card(
        slide,
        Inches(0.55),
        Inches(1.85),
        Inches(4.0),
        Inches(2.85),
        "Why it matters",
        [
            "Protects original creators from silent reuse and counterfeit claims.",
            "Gives moderators stronger evidence before approving or rejecting media.",
            "Builds a path toward safer, more trustworthy visual publishing systems.",
        ],
        RGBColor(232, 240, 254),
    )
    add_text_lines(
        slide,
        Inches(5.05),
        Inches(1.85),
        Inches(3.95),
        Inches(1.1),
        [
            {"text": "SDG 9", "size": 18, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "Industry, Innovation and Infrastructure", "size": 13, "color": DARK},
        ],
        fill=RGBColor(232, 245, 233),
        line=BORDER,
    )
    add_text_lines(
        slide,
        Inches(5.05),
        Inches(3.05),
        Inches(3.95),
        Inches(1.1),
        [
            {"text": "SDG 16", "size": 18, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "Peace, Justice and Strong Institutions", "size": 13, "color": DARK},
        ],
        fill=RGBColor(254, 247, 224),
        line=BORDER,
    )

    # Slide 15
    slide = prs.slides[14]
    add_text_lines(
        slide,
        Inches(1.85),
        Inches(1.65),
        Inches(6.25),
        Inches(2.15),
        [
            {"text": "SecureMedia AI", "size": 26, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER, "space_after": 8},
            {"text": "Protect originality. Prove ownership.", "size": 18, "color": DARK, "align": PP_ALIGN.CENTER, "space_after": 10},
            {"text": "GitHub: github.com/Aadishah17/secure-media", "size": 12.5, "color": BLUE, "align": PP_ALIGN.CENTER, "url": "https://github.com/Aadishah17/secure-media"},
            {"text": "Live app: securemedia-ai-97902534410.us-central1.run.app", "size": 11.5, "color": BLUE, "align": PP_ALIGN.CENTER, "url": SERVICE_URL},
        ],
        fill=WHITE,
        line=BORDER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
